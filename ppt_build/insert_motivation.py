# -*- coding: utf-8 -*-
"""插入「引子」页：为什么是 3 epoch + 50 帧 + 3 seed —— 先测量我们自己的测量误差。
数据来源：早期 Robust-N2N-结果记录.pptx（V3/V4/V5、OOD）+ 本轮多帧/多 seed 实测。
插到第 3 页（阶段1 之前），不触碰用户自制页。
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DECK = r"D:\Desktop\Robust_N2N_实验结果汇报.pptx"
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1D, 0x9E, 0x75); PINK = RGBColor(0xD4, 0x53, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
IN = 914400  # EMU per inch


def run(par, text, size, color, bold=False, italic=False):
    r = par.add_run(); r.text = text
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), F, bold, italic
    r.font.color.rgb = color
    return r


def tbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


prs = Presentation(DECK)
s = prs.slides.add_slide(prs.slide_layouts[0])

# ---- 页眉 ----
run(tbox(s, 0.6*IN, 0.34*IN, 12*IN, 0.32*IN).paragraphs[0],
    "引子 · 先测量我们自己的测量误差", 12, PINK, bold=True)
run(tbox(s, 0.6*IN, 0.66*IN, 12.1*IN, 0.7*IN).paragraphs[0],
    "为什么必须是 3 epoch + 50 帧配对 + 3 seed？", 27, NAVY, bold=True)

# ---- 左：早期记录（单图 · 单次训练）----
run(tbox(s, 0.6*IN, 1.55*IN, 6*IN, 0.3*IN).paragraphs[0],
    "早期记录：单张图 · 单次训练（结果记录 V3~V5 / OOD）", 12.5, NAVY, bold=True)
rows = [("V3", "33.92", "33.25", "−0.67", MUTED),
        ("V4", "33.92", "32.91", "−1.01", MUTED),
        ("V5", "33.92", "34.08", "+0.16", RED),
        ("OOD 5x5x1", "32.22", "31.85", "−0.37", MUTED),
        ("OOD 5x5x5", "33.26", "33.81", "+0.55", MUTED)]
tb = s.shapes.add_table(len(rows) + 1, 4, Emu(int(0.6*IN)), Emu(int(1.95*IN)),
                        Emu(int(6.0*IN)), Emu(int(2.1*IN))).table
for w, c in zip((1.8, 1.4, 1.4, 1.4), range(4)):
    tb.columns[c].width = Emu(int(w*IN))
for j, h in enumerate(["实验", "N2N", "Robust", "差值"]):
    cell = tb.cell(0, j); cell.text = ""
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, h, 10.5, WHITE, bold=True)
for i, (a, b, c, d, col) in enumerate(rows, 1):
    for j, (v, cc, bd) in enumerate([(a, INK, False), (b, INK, False), (c, INK, False), (d, col, True)]):
        cell = tb.cell(i, j); cell.text = ""
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, v, 10.5, cc, bold=bd)
tf = tbox(s, 0.6*IN, 4.2*IN, 6.0*IN, 0.75*IN)
run(tf.paragraphs[0], "「V5 略优于 N2N」这一结论，全部建立在 +0.16 dB 之上。", 11.5, RED, bold=True)

# ---- 右：实测的噪声底 ----
run(tbox(s, 7.0*IN, 1.55*IN, 5.9*IN, 0.3*IN).paragraphs[0],
    "但我们实测了这套测量方法自身的误差", 12.5, NAVY, bold=True)
card = s.shapes.add_shape(5, Emu(int(7.0*IN)), Emu(int(1.95*IN)), Emu(int(5.9*IN)), Emu(int(2.35*IN)))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT
card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
tf = tbox(s, 7.28*IN, 2.15*IN, 5.35*IN, 2.0*IN)
items = [("同一模型、同一场景、不同帧：", "PSNR 28.63 ~ 33.49，std 0.83 dB，极差 4.85 dB"),
         ("同一配置、换一个 seed 重训：", "33.27 / 32.45 / 31.91，摆动 1.36 dB"),
         ("同一配置、同 seed 重跑一次：", "33.951 vs 33.560，差 0.39 dB")]
first = True
for k, v in items:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    run(p, k, 11.5, NAVY, bold=True); run(p, v, 11.5, INK)

# ---- 右下：结论 ----
box = s.shapes.add_shape(5, Emu(int(7.0*IN)), Emu(int(4.5*IN)), Emu(int(5.9*IN)), Emu(int(1.35*IN)))
box.fill.solid(); box.fill.fore_color.rgb = NAVY
box.line.fill.background(); box.shadow.inherit = False
tf = tbox(s, 7.28*IN, 4.72*IN, 5.35*IN, 1.0*IN)
run(tf.paragraphs[0], "+0.16 dB 的「优势」比噪声小一个数量级 —— 那个实验设计没有分辨能力，"
                      "既不能证明方法更好，也不能证明更差。", 12, WHITE, bold=True)

# ---- 底部：对策表 ----
run(tbox(s, 0.6*IN, 5.1*IN, 6.0*IN, 0.3*IN).paragraphs[0], "两种噪声，分别对付", 12.5, NAVY, bold=True)
tf = tbox(s, 0.6*IN, 5.45*IN, 6.0*IN, 1.5*IN)
plan = [("评测噪声（抽到哪张图）→ ", "50 帧平均 + 逐帧配对，标准误压到 ≈0.1 dB"),
        ("训练噪声（seed / 非确定性）→ ", "3 个 seed，判据 mean > std"),
        ("欠训练（1 epoch）→ ", "3 epoch，读数代表模型而非训练暂态")]
first = True
for k, v in plan:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    run(p, k, 11, GREEN, bold=True); run(p, v, 11, INK)

tf = tbox(s, 7.0*IN, 6.05*IN, 5.9*IN, 0.9*IN)
run(tf.paragraphs[0], "验证：新协议当场抓到一个假结论 —— seed42 上 ΔPSNR=+0.762（50/50 帧全赢），"
                      "换 seed 后变成 −0.060 / −0.598，跨 seed 仅 +0.03±0.69（打平）。", 11, INK)

# ---- 移到第 3 页 ----
ids = prs.slides._sldIdLst
new = list(ids)[-1]
ids.remove(new)
ids.insert(2, new)

prs.save(DECK)
print("inserted at position 3; total slides:", len(prs.slides))
