# -*- coding: utf-8 -*-
"""撤掉第8页未经证实的「特征健康度」机制卡片，保留 PSNR 结果，页面改名。
理由：projector 末层是 BatchNorm，而 no-projector 分支的 z 是无 BN 的 native 特征，
两者 std 不可比；且旧监控把 batch 与空间维混在一起，无法排除维度塌缩。
"""
import sys, glob
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

DECK = sys.argv[1] if len(sys.argv) > 1 else r"D:\Desktop\Robust_N2N_实验结果汇报_v2.pptx"
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); RED = RGBColor(0xC0, 0x39, 0x2B)
IN = 914400

prs = Presentation(DECK)
# 找到 projector 消融页（标题含「projector 消融」）
target = None
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and "projector 消融" in sh.text_frame.text:
            target = s; break
    if target: break
if target is None:
    raise SystemExit("未找到 projector 消融页")

# 1) 改标题（kicker + 主标题）
for sh in target.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if "projector 消融" in t and "阶段" in t:
        sh.text_frame.paragraphs[0].runs[0].text = "阶段 5 · projector 消融与待验证机制（e3+bn, w_feat=0.3, 3 seeds）"
    elif "OOD 上挽回约 1 dB" in t:
        sh.text_frame.paragraphs[0].runs[0].text = "projector 显著压平 OOD 衰退（机制待验证）"

# 2) 删掉机制卡片（灰底圆角矩形 + 其文字）
kill = []
for sh in target.shapes:
    txt = sh.text_frame.text if sh.has_text_frame else ""
    if "特征健康度" in txt or "机制证据" in txt:
        kill.append(sh)
    # 卡片本体：y≈4.05in 的圆角矩形
    elif sh.shape_type == 5 and abs(sh.top - int(4.05 * IN)) < int(0.2 * IN):
        kill.append(sh)
for sh in kill:
    sh._element.getparent().remove(sh._element)

# 3) 放一段说明，交代为何撤下
tb = target.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(4.05 * IN)), Emu(int(6.3 * IN)), Emu(int(1.2 * IN)))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
r = p.add_run(); r.text = "机制尚未证实："
r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(11.5), F, True, RED
r = p.add_run()
r.text = ("原先用「ℓ2-归一化特征逐通道 std」作为特征健康度证据，但 projector 末层含 BatchNorm，"
          "而无 projector 分支的 z 是无 BN 的 native 特征 —— 两者 std 不可比；该监控又把 batch 与空间维混算，"
          "无法排除维度塌缩。故撤下该机制结论，仅保留 PSNR 结果。")
r.font.size, r.font.name, r.font.color.rgb = Pt(11), F, INK

p2 = tf.add_paragraph()
r = p2.add_run()
r.text = "下一步以分支隔离的 zero/shuffle 干预 + native 特征的有效秩/协方差谱来检验。"
r.font.size, r.font.name, r.font.italic, r.font.color.rgb = Pt(11), F, True, MUTED

prs.save(DECK)
print("patched:", DECK)
