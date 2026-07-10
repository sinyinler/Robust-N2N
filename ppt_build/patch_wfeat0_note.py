# -*- coding: utf-8 -*-
"""在用户手改过的 pptx 上做手术式修正：澄清 w_feat=0 ≠ N2N 基线。
只改 slide2/slide3 的文字，不触碰用户自制的图片页。
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

DECK = r"D:\Desktop\Robust_N2N_实验结果汇报.pptx"
F = "Microsoft YaHei"
NAVY, INK, MUTED, RED = RGBColor(0x1E, 0x27, 0x61), RGBColor(0x22, 0x22, 0x22), RGBColor(0x6B, 0x6B, 0x66), RGBColor(0xC0, 0x39, 0x2B)


def run(par, text, size, color, bold=False, italic=False):
    r = par.add_run(); r.text = text
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), F, bold, italic
    r.font.color.rgb = color


def add_note(slide, left, top, width, height, text, size=10):
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    run(tf.paragraphs[0], text, size, RED, italic=True)
    return tb


prs = Presentation(DECK)

# ---- slide 2：右卡片下方补一行基线损失说明 ----
s2 = prs.slides[1]
add_note(s2, 6675120, 5850000, 4937760, 620000,
         "注：N2N 基线 = 单向 Charbonnier + 0.01·RTV；Robust 的 w_feat=0 ≠ 此基线"
         "（它仍有 γ=0.1 一致性项、无 RTV、两向重建）。", size=10)

# ---- slide 3：表 B 下方加红色说明 ----
s3 = prs.slides[2]
add_note(s3, 6309360, 4450000, 5486400, 560000,
         "注：w_feat=0 ≠ N2N 基线。它仍是「两向对称重建 + γ=0.1 一致性、无 RTV」；"
         "基线是「单向重建 + RTV 0.01」。", size=10)

# ---- slide 3：重写「初步发现」文本框（shape index 6）----
tf = s3.shapes[6].text_frame
tf.clear(); tf.word_wrap = True
p0 = tf.paragraphs[0]
run(p0, "初步发现（seed42，单 seed，仅供圈定范围）：", 11.5, NAVY, bold=True)
p1 = tf.add_paragraph()
run(p1, "① 分布上 e3+bn 最好，加浅层(e1/e2)全部落进噪声；② 强度需 w_feat≥0.15，太弱(0/0.05/0.1)反而输给基线；", 11.5, INK)
p2 = tf.add_paragraph()
run(p2, "③ 同一 base 下 w_feat 0→0.3 提升 +1.6 dB，这是特征损失的净效果；但 w_feat=0 低于 N2N 的 0.87 dB "
        "主要来自 base 差异（γ 项 / 无 RTV / 两向重建），不能全归因于「缺特征损失」。", 11.5, INK)

prs.save(DECK)
print("patched:", DECK, "| slides:", len(prs.slides))
