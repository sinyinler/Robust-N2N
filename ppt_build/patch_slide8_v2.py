# -*- coding: utf-8 -*-
"""修正版：按页序号精确定位（不再文本匹配，上一版误伤了封面）。
1) 删除封面上误加的说明框；
2) 第 8 页：撤下未证实的「特征健康度」机制卡片 + 结论里「特征不被压缩」的措辞；
   页面改名为「projector 消融与待验证机制」，保留全部 PSNR 结果。
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

DECK = r"D:\Desktop\Robust_N2N_实验结果汇报_v2.pptx"
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); RED = RGBColor(0xC0, 0x39, 0x2B)
IN = 914400


def drop(sh):
    sh._element.getparent().remove(sh._element)


def retext(shape, runs):
    """清空并重建 text_frame，runs = [(text, size, color, bold, italic, newpara), ...]"""
    tf = shape.text_frame
    tf.clear()
    par = tf.paragraphs[0]
    for text, size, color, bold, italic, newpara in runs:
        if newpara:
            par = tf.add_paragraph()
        r = par.add_run(); r.text = text
        r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), F, bold, italic
        r.font.color.rgb = color


prs = Presentation(DECK)

# ---- 1) 封面：删掉上一版误加的说明框 ----
cover = prs.slides[0]
for sh in list(cover.shapes):
    if sh.has_text_frame and sh.text_frame.text.startswith("机制尚未证实"):
        drop(sh); print("已删除封面上误加的说明框")

# ---- 2) 第 8 页 ----
s8 = prs.slides[7]
assert "projector 消融" in s8.shapes[0].text_frame.text, "第 8 页定位失败，请检查页序"

retext(s8.shapes[0], [("阶段 5 · projector 消融与待验证机制（e3+bn, w_feat=0.3, 3 seeds）",
                       12, RGBColor(0xD4, 0x53, 0x7E), True, False, False)])
retext(s8.shapes[1], [("projector 显著压平 OOD 衰退（机制待验证）", 27, NAVY, True, False, False)])

# 删机制卡片：shape 4 = 灰底圆角矩形，shape 5 = 卡片文字
for sh in list(s8.shapes):
    txt = sh.text_frame.text if sh.has_text_frame else ""
    if "机制证据" in txt or "特征健康度" in txt:
        drop(sh)
    elif sh.shape_type == 5 and abs(sh.top - int(4.05 * IN)) < int(0.15 * IN):
        drop(sh)

# 结论：去掉「特征不被压缩」这一未证实断言
concl = None
for sh in s8.shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("结论："):
        concl = sh; break
if concl is not None:
    retext(concl, [
        ("结论：", 11.5, NAVY, True, False, False),
        ("projector 有大幅、可测量、三 seed 可复现的作用（OOD 挽回约 1 dB、跨 seed 方差减半），并非装饰。",
         11.5, INK, False, False, False),
        ("但仍未赢 N2N，且机制未证实：", 11.5, RED, True, False, True),
        ("最 OOD 层为 −0.67±0.46（|mean|>std，是稳定的输）。projector 把「崩盘」变成「小输」，治标未治本。",
         11.5, INK, False, False, False),
    ])

# 在腾出的位置补一段「为何撤下机制结论」
tb = s8.shapes.add_textbox(Emu(int(0.6 * IN)), Emu(int(4.05 * IN)), Emu(int(6.3 * IN)), Emu(int(1.25 * IN)))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
for text, size, col, bold, ital in [
    ("机制尚未证实：", 11.5, RED, True, False),
    ("原先以「ℓ2-归一化特征逐通道 std」作为特征健康度证据。但 projector 末层含 BatchNorm，"
     "而无 projector 分支的 z 是无 BN 的 native 特征 —— 两者 std 不可比；该监控又把 batch 与空间维混算，"
     "无法排除维度塌缩。故撤下此机制结论，仅保留 PSNR 结果。", 11, INK, False, False),
]:
    r = p.add_run(); r.text = text
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), F, bold, ital
    r.font.color.rgb = col
p2 = tf.add_paragraph()
r = p2.add_run()
r.text = "下一步：分支隔离的 zero/shuffle 干预 + native 特征的有效秩 / 协方差谱来检验。"
r.font.size, r.font.name, r.font.italic, r.font.color.rgb = Pt(11), F, True, MUTED

prs.save(DECK)
print("patched slide 8; total slides:", len(prs.slides))
