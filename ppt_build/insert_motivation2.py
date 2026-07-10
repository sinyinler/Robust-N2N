# -*- coding: utf-8 -*-
"""替换「引子」页：证据改用早期 1-epoch + 单图的五张权重扫描表本身。
核心自证：同一配置分属两张表、两次独立训练 → 33.951 vs 33.560，差 0.391 dB。
"""
import sys
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r"D:\Desktop\Robust_N2N_实验结果汇报_v2.pptx"
DST = sys.argv[1] if len(sys.argv) > 1 else SRC
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1D, 0x9E, 0x75); PINK = RGBColor(0xD4, 0x53, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
IN = 914400


def run(par, text, size, color, bold=False, italic=False):
    r = par.add_run(); r.text = text
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(size), F, bold, italic
    r.font.color.rgb = color
    return r


def tbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


prs = Presentation(SRC)

# ---- 删掉旧的引子页（当前第 3 页）----
ids = prs.slides._sldIdLst
ids.remove(list(ids)[2])

s = prs.slides.add_slide(prs.slide_layouts[0])
run(tbox(s, 0.6*IN, 0.34*IN, 12*IN, 0.32*IN).paragraphs[0],
    "引子 · 早期调参：我们其实在噪声里挑「最优」", 12, PINK, bold=True)
run(tbox(s, 0.6*IN, 0.66*IN, 12.1*IN, 0.7*IN).paragraphs[0],
    "为什么必须是 3 epoch + 50 帧配对 + 3 seed？", 27, NAVY, bold=True)

# ===== 左上：证据一（同一配置两次训练）=====
run(tbox(s, 0.6*IN, 1.5*IN, 6.3*IN, 0.3*IN).paragraphs[0],
    "证据一：同一配置，两次独立训练，结果差 0.391 dB", 12.5, NAVY, bold=True)
rows = [("encoder1 扫描表（e1=0.7 行）", "33.951", INK, False),
        ("γ 扫描表（γ=0.1 行）", "33.560", INK, False),
        ("差值（同一配置！）", "0.391 dB", RED, True)]
tb = s.shapes.add_table(4, 2, Emu(int(0.6*IN)), Emu(int(1.85*IN)),
                        Emu(int(6.3*IN)), Emu(int(1.55*IN))).table
tb.columns[0].width = Emu(int(4.3*IN)); tb.columns[1].width = Emu(int(2.0*IN))
for j, h in enumerate(["数据来源", "PSNR"]):
    c = tb.cell(0, j); c.text = ""
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, h, 10.5, WHITE, bold=True)
for i, (a, b, col, bd) in enumerate(rows, 1):
    for j, (v, cc) in enumerate([(a, INK if not bd else RED), (b, col)]):
        c = tb.cell(i, j); c.text = ""
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, v, 10.5, cc, bold=bd)
tf = tbox(s, 0.6*IN, 3.5*IN, 6.3*IN, 0.5*IN)
run(tf.paragraphs[0], "两行的 γ=0.1, w_feat=0.05, bn=1, e3=1, e2=0.6, e1=0.7 完全一致 —— "
                      "唯一差别是「跑了两次」。", 10.5, MUTED, italic=True)

# ===== 左下：证据二（扫描是抖动不是趋势）=====
run(tbox(s, 0.6*IN, 4.1*IN, 6.3*IN, 0.3*IN).paragraphs[0],
    "证据二：扫描曲线是抖动，不是趋势", 12.5, NAVY, bold=True)
tf = tbox(s, 0.6*IN, 4.5*IN, 6.3*IN, 2.0*IN)
ev = [("w_feat 扫描：", "0.04→33.637 ｜ 0.05→33.814（当时选它）｜ 0.06→33.651 —— 相邻点上下跳 0.16~0.18"),
      ("encoder2 扫描：", "选了 0.6→33.776，但 0.9→33.868 其实更高"),
      ("encoder1 扫描：", "0.7→33.951 与 0.8→33.953，相差 0.002，根本分不开"),
      ("整条扫描线的起伏 0.3~0.5 dB，", "与「同配置重跑」的 0.391 dB 同量级 → 挑的是噪声")]
first = True
for k, v in ev:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    run(p, "• " + k, 11, RED if k.startswith("整条") else NAVY, bold=True)
    run(p, v, 11, INK)

# ===== 右上：结论 =====
box = s.shapes.add_shape(5, Emu(int(7.15*IN)), Emu(int(1.85*IN)), Emu(int(5.6*IN)), Emu(int(1.75*IN)))
box.fill.solid(); box.fill.fore_color.rgb = NAVY
box.line.fill.background(); box.shadow.inherit = False
tf = tbox(s, 7.42*IN, 2.05*IN, 5.05*IN, 1.4*IN)
run(tf.paragraphs[0], "这五张表都是「1 epoch + 单张图」的读数。当测量噪声（0.39 dB）"
                      "和权重带来的差异（0.3~0.5 dB）一样大时，逐个权重贪心挑最优，"
                      "挑出来的是噪声尖峰，不是真实最优。", 12, WHITE, bold=True)

# ===== 右中：对策 =====
run(tbox(s, 7.15*IN, 3.8*IN, 5.6*IN, 0.3*IN).paragraphs[0], "两种噪声，分别对付", 12.5, NAVY, bold=True)
tf = tbox(s, 7.15*IN, 4.2*IN, 5.6*IN, 1.5*IN)
plan = [("评测噪声（抽到哪张图）→ ", "50 帧平均 + 逐帧配对，标准误压到 ≈0.1 dB"),
        ("训练噪声（seed / 非确定性）→ ", "3 个 seed，判据 mean > std"),
        ("欠训练（1 epoch）→ ", "3 epoch，读数代表模型而非训练暂态")]
first = True
for k, v in plan:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    run(p, k, 11, GREEN, bold=True); run(p, v, 11, INK)

# ===== 右下：验证 =====
card = s.shapes.add_shape(5, Emu(int(7.15*IN)), Emu(int(5.5*IN)), Emu(int(5.6*IN)), Emu(int(1.4*IN)))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT
card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
tf = tbox(s, 7.42*IN, 5.68*IN, 5.05*IN, 1.05*IN)
p = tf.paragraphs[0]
run(p, "新协议当场抓到一个假结论：", 11, RED, bold=True)
run(p, "seed42 上 ΔPSNR = +0.762（50/50 帧全赢）；换两个 seed 重训变成 "
       "−0.060 / −0.598，跨 seed 仅 +0.03±0.69 —— 打平。", 11, INK)

# ---- 移到第 3 页 ----
new = list(ids)[-1]
ids.remove(new); ids.insert(2, new)
prs.save(DST)
print("saved:", DST, "| slides:", len(prs.slides))
