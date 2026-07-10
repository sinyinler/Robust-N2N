# -*- coding: utf-8 -*-
"""追加：谢谢页 + 附录（按证据顺序放图）。"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

DECK = r"D:\Desktop\Robust_N2N_实验结果汇报_v2.pptx"
FIGS = os.path.join(os.path.dirname(__file__), "figs")
DESK = r"D:\Desktop"
ABL = r"D:\Desktop\ablation_feat"
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); PINK = RGBColor(0xD4, 0x53, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
IN = 914400
W, H = 13.333, 7.5


def run(par, t, sz, col, bold=False, italic=False):
    r = par.add_run(); r.text = t
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(sz), F, bold, italic
    r.font.color.rgb = col; return r


def tbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Emu(int(l*IN)), Emu(int(t*IN)), Emu(int(w*IN)), Emu(int(h*IN)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


prs = Presentation(DECK)

# ================= 谢谢页 =================
s = prs.slides.add_slide(prs.slide_layouts[0])
bg = s.shapes.add_shape(1, 0, 0, Emu(int(W*IN)), Emu(int(H*IN)))
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
bg.line.fill.background(); bg.shadow.inherit = False
tf = tbox(s, 0.9, 3.0, 11.5, 1.2)
run(tf.paragraphs[0], "谢谢！", 60, WHITE, bold=True)
tf = tbox(s, 0.95, 4.35, 11.5, 0.5)
run(tf.paragraphs[0], "Robust-N2N · 跨视图特征一致性 · 实验结果汇报", 16, RGBColor(0xCA, 0xDC, 0xFC))

# ================= 附录扉页 =================
s = prs.slides.add_slide(prs.slide_layouts[0])
run(tbox(s, 0.6, 0.34, 12, 0.32).paragraphs[0], "APPENDIX", 12, PINK, bold=True)
run(tbox(s, 0.6, 0.66, 12.1, 0.8).paragraphs[0], "附录 · 全部证据图（按实验顺序）", 30, NAVY, bold=True)
tf = tbox(s, 0.6, 1.9, 12.1, 4.5)
items = [
    ("A1", "引子：噪声底 —— 同一配置两次训练差 0.391 dB；扫描曲线是抖动"),
    ("A2", "阶段1：11 条臂的多帧 PSNR 概览（mean ± std）"),
    ("A3", "阶段1：配对 ΔPSNR（绿=稳定赢，灰=噪声内）"),
    ("A4", "阶段1：11 条臂 + 基线的逐帧 PSNR 曲线"),
    ("A5", "阶段2：多 seed 确认 —— seed42 的「冠军」是运气"),
    ("A6", "阶段3：OOD 泛化测试 —— 越 OOD 越差（三 seed 一致、单调）"),
    ("A7", "阶段5：projector 消融 —— OOD 上挽回约 1 dB"),
    ("A8", "阶段6：最终干净消融 —— 唯一差异=特征损失项"),
    ("A9~A11", "阶段6：三个 seed 的逐帧配对曲线（N2N+特征损失 vs N2N）"),
]
first = True
for k, v in items:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    run(p, f"{k}　", 14, PINK, bold=True); run(p, v, 14, INK)

# ================= 图页 =================
figs = [
    (os.path.join(FIGS, "A1_引子_噪声底.png"), "A1 · 引子：噪声与信号同量级",
     "左：w_feat 扫描相邻点上下跳 0.16~0.18 dB。右：encoder1 扫描表与 γ 扫描表里配置完全相同的两行，读数差 0.391 dB。"),
    (os.path.join(DESK, "overview_psnr.png"), "A2 · 阶段1：各臂多帧 PSNR 概览",
     "level4、seed42、3 epoch。误差棒为 50 帧的 std；蓝色虚线带为 N2N 基线 ± std。"),
    (os.path.join(DESK, "dpsnr.png"), "A3 · 阶段1：配对 ΔPSNR",
     "绿柱 = mean > std（稳定赢过 N2N）；灰柱 = 差异落在噪声内。w_feat=0 一条最差。"),
    (os.path.join(DESK, "curves_psnr.png"), "A4 · 阶段1：逐帧 PSNR 曲线",
     "11 条臂 + N2N 基线（黑）在同 50 帧上的逐帧读数，可见帧间波动远大于臂间差异。"),
    (os.path.join(FIGS, "A5_阶段2_多seed.png"), "A5 · 阶段2：多 seed 确认",
     "同一配置换 seed：w_feat=0.3 从 +0.762 变 −0.060 / −0.598；w_feat=0.15 在 seed43 崩 −2.374。"),
    (os.path.join(FIGS, "A6_阶段3_OOD趋势.png"), "A6 · 阶段3：OOD 泛化测试（无 projector）",
     "训练于 level4。ΔPSNR 随 OOD 加深单调下降，三 seed 一致；level1 均值 −1.63 dB。"),
    (os.path.join(FIGS, "A7_阶段5_projector消融.png"), "A7 · 阶段5：projector 消融",
     "带 projector 后 OOD 衰退被大幅压平（level1 −1.63 → −0.67）。改善随 OOD 加深而增大。"),
    (os.path.join(FIGS, "A8_阶段6_干净消融.png"), "A8 · 阶段6：最终干净消融",
     "同 seed / 同 batch / 同损失，唯一差异 = 特征损失项。跨 seed −0.362 ± 0.479。"),
    (os.path.join(ABL, "s42", "psnr_curve.png"), "A9 · 阶段6：seed42 逐帧配对曲线",
     "N2N + 特征损失 vs 纯 N2N，同 50 帧。配对 ΔPSNR = +0.109 ± 0.400（28/50，噪声内）。"),
    (os.path.join(ABL, "s43", "psnr_curve.png"), "A10 · 阶段6：seed43 逐帧配对曲线",
     "配对 ΔPSNR = −0.349 ± 0.295（3/50，稳定更差）。"),
    (os.path.join(ABL, "s44", "psnr_curve.png"), "A11 · 阶段6：seed44 逐帧配对曲线",
     "配对 ΔPSNR = −0.848 ± 0.690（8/50，稳定更差）。"),
]

added = 0
for path, title, cap in figs:
    if not os.path.exists(path):
        print("[SKIP 缺图]", path); continue
    s = prs.slides.add_slide(prs.slide_layouts[0])
    run(tbox(s, 0.6, 0.34, 12, 0.5).paragraphs[0], title, 20, NAVY, bold=True)
    # 等比缩放，居中放入 (0.6,1.05) ~ (12.73, 5.95)
    iw, ih = Image.open(path).size
    box_w, box_h = 12.13, 4.9
    sc = min(box_w / iw, box_h / ih)
    w_in, h_in = iw * sc, ih * sc
    s.shapes.add_picture(path, Emu(int((W - w_in) / 2 * IN)), Emu(int((1.05 + (box_h - h_in) / 2) * IN)),
                         Emu(int(w_in * IN)), Emu(int(h_in * IN)))
    run(tbox(s, 0.6, 6.15, 12.13, 0.9).paragraphs[0], cap, 12, MUTED)
    added += 1

prs.save(DECK)
print(f"added {added} figure slides; total slides = {len(prs.slides)}")
