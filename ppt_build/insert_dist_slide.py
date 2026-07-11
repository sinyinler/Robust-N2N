# -*- coding: utf-8 -*-
"""追加一页：分布扫描（128 projector + 池化 + 3 seed）—— 分布轴不可分辨。"""
import os, sys, glob, shutil
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

DECK = r"D:\Desktop\Robust_N2N_实验结果汇报_v2.pptx"
FIG = os.path.join(os.path.dirname(__file__), "figs", "A9_分布扫描_不可分辨.png")
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); PINK = RGBColor(0xD4, 0x53, 0x7E); RED = RGBColor(0xC0, 0x39, 0x2B)
IN = 914400; W, H = 13.333, 7.5

# 备份
bak = DECK.replace(".pptx", f".bak3-{__import__('time').strftime('%H%M%S')}.pptx")
shutil.copy(DECK, bak)

prs = Presentation(DECK)
# 找到「谢谢」页的位置，插在它前面（正文最后一页之后）
titles = []
for sl in prs.slides:
    t = ""
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.strip().split("\n")[0]; break
    titles.append(t)
thanks_idx = next((i for i, t in enumerate(titles) if t.startswith("谢谢")), len(prs.slides))


def tbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Emu(int(l*IN)), Emu(int(t*IN)), Emu(int(w*IN)), Emu(int(h*IN)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def run(par, txt, sz, col, bold=False):
    r = par.add_run(); r.text = txt
    r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(sz), F, bold, col
    return r


s = prs.slides.add_slide(prs.slide_layouts[0])
run(tbox(s, 0.6, 0.34, 12, 0.32).paragraphs[0], "阶段 7 · 128 projector + 区域级一致性 + 分布扫描（3 seed）", 12, PINK, bold=True)
run(tbox(s, 0.6, 0.66, 12.1, 0.7).paragraphs[0], "换了 projector、换了池化、扫了全部分布 —— 结论不变：分布轴不可分辨", 24, NAVY, bold=True)

iw, ih = Image.open(FIG).size
bw, bh = 12.2, 4.35
sc = min(bw / iw, bh / ih); w_in, h_in = iw * sc, ih * sc
s.shapes.add_picture(FIG, Emu(int((W - w_in) / 2 * IN)), Emu(int(1.5 * IN)), Emu(int(w_in * IN)), Emu(int(h_in * IN)))

tf = tbox(s, 0.6, 6.05, 12.1, 1.3)
p = tf.paragraphs[0]
run(p, "① 分布不可分辨：", 12.5, RED, bold=True)
run(p, "6 条臂跨 3 seed 均值全部落在 0 附近，误差棒无一不跨过 0；最优臂「深重浅轻」在 seed2413 排第 1、seed187 排第 6。", 12.5, INK)
p = tf.add_paragraph()
run(p, "② feat 与 PSNR 无关联：", 12.5, RED, bold=True)
run(p, "损失三 seed 都轻松拉到 −0.95（基线仅 ≈0），但对应 PSNR 一 seed 全正、两 seed 全负 —— 对齐程度和去噪质量之间没有稳定关系。", 12.5, INK)
p = tf.add_paragraph()
run(p, "结论：", 12.5, NAVY, bold=True)
run(p, "128 维 projector、区域级一致性、六种分布都试过 —— 同级 noisy-pair 的跨视图一致性，与去噪质量无可测、可复现的关联。", 12.5, INK)

# 移到「谢谢」前
ids = prs.slides._sldIdLst
new = list(ids)[-1]; ids.remove(new); ids.insert(thanks_idx, new)
prs.save(DECK)
print(f"inserted before 谢谢(idx {thanks_idx}); total {len(prs.slides)}; backup {os.path.basename(bak)}")
