# -*- coding: utf-8 -*-
"""追加「最终干净消融」页：同 seed / 同 batch / 同损失，唯一差异 = 特征损失项。"""
import sys, glob
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r"D:\Desktop\Robust_N2N_实验结果汇报_v2.pptx"
DST = sys.argv[1] if len(sys.argv) > 1 else SRC
F = "Microsoft YaHei"
NAVY = RGBColor(0x1E, 0x27, 0x61); INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x6B, 0x66); RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1D, 0x9E, 0x75); PINK = RGBColor(0xD4, 0x53, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
LINE = RGBColor(0xD8, 0xDC, 0xE3)
IN = 914400


def run(par, t, sz, col, bold=False, italic=False):
    r = par.add_run(); r.text = t
    r.font.size, r.font.name, r.font.bold, r.font.italic = Pt(sz), F, bold, italic
    r.font.color.rgb = col; return r


def tbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def table(s, data, l, t, w, h, colw, fs=10.5):
    tb = s.shapes.add_table(len(data), len(data[0]), Emu(int(l)), Emu(int(t)),
                            Emu(int(w)), Emu(int(h))).table
    for i, cw in enumerate(colw):
        tb.columns[i].width = Emu(int(cw * IN))
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            txt, col, bd = cell if isinstance(cell, tuple) else (cell, INK, False)
            c = tb.cell(i, j); c.text = ""
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run(p, txt, fs, WHITE if i == 0 else col, bold=True if i == 0 else bd)
    return tb


prs = Presentation(SRC)
s = prs.slides.add_slide(prs.slide_layouts[0])

run(tbox(s, 0.6*IN, 0.34*IN, 12*IN, 0.32*IN).paragraphs[0],
    "阶段 6 · 最终干净消融（同 seed / 同 batch / 唯一差异 = 特征损失项）", 12, PINK, bold=True)
run(tbox(s, 0.6*IN, 0.66*IN, 12.1*IN, 0.7*IN).paragraphs[0],
    "把所有混淆去掉后：特征损失没有提升，结构指标一致略差", 27, NAVY, bold=True)

# ---- 左：主表 ----
run(tbox(s, 0.6*IN, 1.5*IN, 7.2*IN, 0.3*IN).paragraphs[0],
    "L = Charb(f(x_in), x_tgt) + 0.01·RTV(f(x_in))  ［+ w_feat·L_feat］    level4 · 3 epoch · batch 24",
    11, MUTED, italic=True)
data = [["seed", "N2N 基线", "N2N + 特征损失", "配对 ΔPSNR", "赢/50", "判定"],
        ["42", "32.434±1.182", "32.543±0.848", ("+0.109±0.400", MUTED, True), "28/50", ("噪声内", MUTED, False)],
        ["43", "32.556±1.042", "32.207±0.752", ("−0.349±0.295", RED, True), "3/50", ("稳定更差", RED, True)],
        ["44", "32.907±0.816", "32.059±1.202", ("−0.848±0.690", RED, True), "8/50", ("稳定更差", RED, True)],
        [("跨 seed", NAVY, True), ("32.632±0.245", NAVY, True), ("32.270±0.248", NAVY, True),
         ("−0.362±0.479", RED, True), "—", ("打平·点估计为负", RED, True)]]
table(s, data, 0.6*IN, 1.9*IN, 7.2*IN, 2.2*IN, [0.65, 1.45, 1.5, 1.4, 0.7, 1.25])

# ---- 左下：结构指标 ----
run(tbox(s, 0.6*IN, 4.3*IN, 7.2*IN, 0.3*IN).paragraphs[0],
    "结构指标：三个 seed、两个指标，方向完全一致（全为负）", 12, NAVY, bold=True)
table(s, [["", "seed42", "seed43", "seed44"],
          [("ΔMSSIM", INK, True), ("−0.0049", RED, False), ("−0.0148", RED, False), ("−0.0056", RED, False)],
          [("Δr", INK, True), ("−0.0048", RED, False), ("−0.0092", RED, False), ("−0.0071", RED, False)]],
      0.6*IN, 4.65*IN, 7.2*IN, 0.95*IN, [1.5, 1.9, 1.9, 1.9])
run(tbox(s, 0.6*IN, 5.75*IN, 7.2*IN, 0.4*IN).paragraphs[0],
    "PSNR 说「打平偏负」，MSSIM / r 说「一致地略差」—— 这不是噪声。", 11.5, RED, bold=True)

# ---- 右上：结论 ----
box = s.shapes.add_shape(5, Emu(int(8.05*IN)), Emu(int(1.85*IN)), Emu(int(4.7*IN)), Emu(int(1.85*IN)))
box.fill.solid(); box.fill.fore_color.rgb = NAVY
box.line.fill.background(); box.shadow.inherit = False
tf = tbox(s, 8.3*IN, 2.05*IN, 4.2*IN, 1.5*IN)
run(tf.paragraphs[0], "一旦把 γ 一致性项、RTV、两向重建、batch、seed 全部对齐，"
                      "早期看到的正向增益就消失了。特征损失在 in-distribution 上不成立。", 12, WHITE, bold=True)

# ---- 右中：附带发现 ----
run(tbox(s, 8.05*IN, 3.95*IN, 4.7*IN, 0.3*IN).paragraphs[0], "两个附带发现", 12.5, NAVY, bold=True)
card = s.shapes.add_shape(5, Emu(int(8.05*IN)), Emu(int(4.3*IN)), Emu(int(4.7*IN)), Emu(int(2.4*IN)))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT
card.line.color.rgb = LINE; card.line.width = Pt(1); card.shadow.inherit = False
tf = tbox(s, 8.3*IN, 4.5*IN, 4.2*IN, 2.05*IN)
p = tf.paragraphs[0]
run(p, "① 训练不稳定不是特征损失的锅。", 11.5, GREEN, bold=True)
p = tf.add_paragraph()
run(p, "去掉 γ 项与两向重建后，跨 seed std 从 ±0.685 降到 ±0.248 —— 与纯 N2N（±0.245）持平。"
       "先前的剧烈摆动来自输出级 γ 项 + 无 projector 的特征压缩。", 11, INK)
p = tf.add_paragraph()
run(p, "② batch 混淆可忽略。", 11.5, GREEN, bold=True)
p = tf.add_paragraph()
run(p, "batch24 基线 32.434 vs batch48 32.509，仅差 0.075 dB —— 先前所有结论无需因 batch 修正。", 11, INK)

prs.save(DST)
print("saved:", DST, "| slides:", len(prs.slides))
